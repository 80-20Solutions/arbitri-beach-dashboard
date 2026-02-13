from pptx import Presentation
import copy

prs = Presentation('Conferenza_v2.pptx')
print(f"Slide originali: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            title = shape.text_frame.text[:60]
            break
    print(f"  Slide {i+1}: {title}")

# Remove slides 1-4 (indices 0-3)
# python-pptx: need to remove from XML
slide_ids = list(prs.slides._sldIdLst)
for i in range(4):
    rId = slide_ids[i].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(slide_ids[i])

prs.save('Conferenza_v3.pptx')
print(f"\nSlide rimanenti: {len(prs.slides)}")
print("Salvato come Conferenza_v3.pptx")
