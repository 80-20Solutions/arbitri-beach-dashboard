from pptx import Presentation
prs = Presentation(r'C:\Users\KreshOS\Desktop\00-arbitr\00-arbitri\Conferenza.pptx')
print(f'Slide: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    print(f'\n--- Slide {i+1} (layout: {slide.slide_layout.name}) ---')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t: print(f'  {t}')
        if shape.has_table:
            print('  [TABELLA]')
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                print('    ' + ' | '.join(cells))
        if shape.shape_type == 13:  # picture
            print(f'  [IMMAGINE: {shape.width}x{shape.height}]')
