from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree
import io

prs = Presentation('Conferenza.pptx')
slides = list(prs.slides)

# Get the blank layout (all slides use BLANK)
blank_layout = prs.slide_layouts[6]  # typically BLANK

# Helper: clone base shapes (background, freeform, logo) from slide 5 (index 4)
# Slides 5,6,7 already have the right base shapes. We'll modify them in place
# and add new slides 8,9.

def get_base_xml(slide):
    """Get XML elements for: AutoShape background, Freeform, logo picture"""
    elements = []
    for shape in slide.shapes:
        if shape.name == 'Google Shape;84;p1' or shape.name == 'Google Shape;85;p1' or shape.name == 'image7.png':
            elements.append(deepcopy(shape._element))
    return elements

def set_title(slide, text):
    """Find the title textbox and set its text"""
    for shape in slide.shapes:
        if shape.name == 'Google Shape;86;p1' and hasattr(shape, 'text_frame'):
            if 'Titolo presentazione' in shape.text or shape.text.strip() == '':
                # Reuse existing formatting
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.text = text
                        return
                # If no runs, add one
                p = shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.name = 'Aharoni'
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
                return

def add_textbox(slide, left, top, width, height, text, font_name='Calibri', font_size=Pt(11), bold=False, color=RGBColor(0,0,0), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_subtitle(slide, left, top, width, height, text):
    return add_textbox(slide, left, top, width, height, text,
                       font_name='Aharoni', font_size=Pt(18), bold=True,
                       color=RGBColor(0xC0, 0x00, 0x00))

def create_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_header_cell(cell, text, font_size=Pt(9)):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'Roboto'
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Set cell fill
    tc = cell._tc
    tcPr = tc.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
    if tcPr is None:
        nsmap = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        tcPr = etree.SubElement(tc, f'{nsmap}tcPr')
    fill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgb = etree.SubElement(fill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgb.set('val', '356854')

def style_data_cell(cell, text, font_size=Pt(9), bold=False, align=PP_ALIGN.CENTER, fill_color=None):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Roboto'
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    if fill_color:
        tc = cell._tc
        tcPr = tc.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
        if tcPr is None:
            tcPr = etree.SubElement(tc, '{http://schemas.openxmlformats.org/drawingml/2006/main}tcPr')
        fill = etree.SubElement(tcPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        srgb = etree.SubElement(fill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgb.set('val', fill_color)

def add_new_slide_with_base(prs):
    """Add a new blank slide and copy base shapes from slide 1"""
    slide = prs.slides.add_slide(blank_layout)
    # Copy base elements from slide 1
    src = prs.slides[0]
    for shape in src.shapes:
        if shape.name in ('Google Shape;84;p1', 'Google Shape;85;p1', 'image7.png'):
            el = deepcopy(shape._element)
            slide.shapes._spTree.append(el)
    # Add title textbox (same position as slides 2-7)
    txBox = slide.shapes.add_textbox(Emu(2518913), Emu(103516), Emu(7375500), Emu(954300))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Titolo presentazione'
    run.font.name = 'Aharoni'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
    # Name it so set_title can find it
    txBox.name = 'Google Shape;86;p1'
    return slide

def add_bullet_textbox(slide, left, top, width, height, items, font_size=Pt(11), color=RGBColor(0x33,0x33,0x33)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f'• {item}'
        run.font.name = 'Calibri'
        run.font.size = font_size
        run.font.color.rgb = color
    return txBox

# ============================================
# SLIDE 5: Arbitri Beach Volley - Panoramica
# ============================================
slide5 = prs.slides[4]
set_title(slide5, 'Arbitri Beach Volley - Panoramica')

# Key numbers textbox
add_textbox(slide5, Emu(410714), Emu(1085649), Emu(3043687), Emu(500000),
            '65 Arbitri  |  15 Supervisori  |  Età media: 40 anni',
            font_name='Calibri', font_size=Pt(14), bold=True, color=RGBColor(0x33,0x33,0x33))

add_textbox(slide5, Emu(410714), Emu(1500000), Emu(3043687), Emu(350000),
            'Arbitri attivi 2024-25: 25 su 65 (38%)',
            font_name='Calibri', font_size=Pt(12), bold=False, color=RGBColor(0x33,0x33,0x33))

# Table: distribution by committee
comitati = [
    ('Romagna Uno', '18'),
    ('Bologna', '14'),
    ('Modena', '10'),
    ('Parma', '5'),
    ('Ravenna', '5'),
    ('Piacenza', '3'),
    ('Reggio Emilia', '3'),
    ('Ferrara', '2'),
    ('TOTALE', '65'),
]
tbl = create_table(slide5, Emu(3500000), Emu(1050000), Emu(4200000), Emu(2300000), 10, 2)
style_header_cell(tbl.cell(0,0), 'Comitato')
style_header_cell(tbl.cell(0,1), 'Nr. Arbitri')
for i, (ct, nr) in enumerate(comitati):
    bold = (ct == 'TOTALE')
    style_data_cell(tbl.cell(i+1, 0), ct, bold=bold, align=PP_ALIGN.LEFT)
    style_data_cell(tbl.cell(i+1, 1), nr, bold=bold)

# Qualifiche section
add_subtitle(slide5, Emu(410714), Emu(3400000), Emu(3043687), Emu(400000), 'Distribuzione per Qualifica')

qual_tbl = create_table(slide5, Emu(3500000), Emu(3500000), Emu(4200000), Emu(1100000), 4, 2)
style_header_cell(qual_tbl.cell(0,0), 'Qualifica')
style_header_cell(qual_tbl.cell(0,1), 'Nr. Arbitri')
for i, (q, n) in enumerate([('PR - Provinciale', '~20'), ('RE - Regionale', '~32'), ('NA - Nazionale', '~13')]):
    style_data_cell(qual_tbl.cell(i+1, 0), q, align=PP_ALIGN.LEFT)
    style_data_cell(qual_tbl.cell(i+1, 1), n)

# Age distribution
add_subtitle(slide5, Emu(8000000), Emu(1050000), Emu(3500000), Emu(400000), 'Fasce di Età')

age_tbl = create_table(slide5, Emu(8000000), Emu(1500000), Emu(3500000), Emu(1600000), 6, 2)
style_header_cell(age_tbl.cell(0,0), 'Fascia')
style_header_cell(age_tbl.cell(0,1), 'Nr. Arbitri')
for i, (f, n) in enumerate([('<20', '5'), ('20-30', '15'), ('30-40', '11'), ('40-50', '15'), ('>50', '19')]):
    style_data_cell(age_tbl.cell(i+1, 0), f)
    style_data_cell(age_tbl.cell(i+1, 1), n)


# ============================================
# SLIDE 6: Attività Beach Volley
# ============================================
slide6 = prs.slides[5]
set_title(slide6, 'Attività Beach Volley')

# Activity table
add_subtitle(slide6, Emu(410714), Emu(1085649), Emu(3043687), Emu(400000), 'Tornei per Stagione')

act_data = [('2018','1'),('2019','36'),('2020','3'),('2021','8'),('2022','8'),('2023','14'),('2024','16'),('2025','6*')]
act_tbl = create_table(slide6, Emu(410714), Emu(1500000), Emu(3500000), Emu(2400000), 9, 2)
style_header_cell(act_tbl.cell(0,0), 'Stagione')
style_header_cell(act_tbl.cell(0,1), 'Nr. Tornei')
for i, (y, t) in enumerate(act_data):
    bold = (y == '2019')
    style_data_cell(act_tbl.cell(i+1, 0), y, bold=bold)
    style_data_cell(act_tbl.cell(i+1, 1), t, bold=bold)

add_textbox(slide6, Emu(410714), Emu(3950000), Emu(3500000), Emu(300000),
            '* stagione in corso | Record: 36 tornei nel 2019 (pre-COVID)',
            font_name='Calibri', font_size=Pt(9), bold=False, color=RGBColor(0x66,0x66,0x66))

# Top arbitri
add_subtitle(slide6, Emu(4300000), Emu(1085649), Emu(4000000), Emu(400000), 'Top Arbitri per Designazioni')

top_data = [
    ('Marchetti Marco', '52', 'Romagna Uno', '37'),
    ('Rosati Andrea', '52', 'Romagna Uno', '32'),
    ('Tramontano Ciro', '30', 'Romagna Uno', '29'),
    ('Brunelli Andrea', '61', 'Ravenna', '26'),
    ('Drei Marco', '62', 'Bologna', '26'),
    ('Greco Luigi', '46', 'Romagna Uno', '25'),
]
top_tbl = create_table(slide6, Emu(4300000), Emu(1500000), Emu(7200000), Emu(1800000), 7, 4)
for ci, h in enumerate(['Arbitro', 'Età', 'Comitato', 'Tornei']):
    style_header_cell(top_tbl.cell(0, ci), h)
for i, (name, age, ct, tornei) in enumerate(top_data):
    style_data_cell(top_tbl.cell(i+1, 0), name, align=PP_ALIGN.LEFT)
    style_data_cell(top_tbl.cell(i+1, 1), age)
    style_data_cell(top_tbl.cell(i+1, 2), ct)
    style_data_cell(top_tbl.cell(i+1, 3), tornei, bold=True)

# Candidati promozione
add_subtitle(slide6, Emu(4300000), Emu(3500000), Emu(5000000), Emu(400000), 'Candidati Promozione (≥10 tornei & <30 anni)')

cand_tbl = create_table(slide6, Emu(4300000), Emu(3950000), Emu(5500000), Emu(800000), 3, 3)
for ci, h in enumerate(['Arbitro', 'Età', 'Tornei']):
    style_header_cell(cand_tbl.cell(0, ci), h)
for i, (name, age, t) in enumerate([('Avanzolini Alberto', '27', '15'), ('Carota Giacomo', '19', '10')]):
    style_data_cell(cand_tbl.cell(i+1, 0), name, align=PP_ALIGN.LEFT)
    style_data_cell(cand_tbl.cell(i+1, 1), age)
    style_data_cell(cand_tbl.cell(i+1, 2), t, bold=True)


# ============================================
# SLIDE 7: Criticità e Opportunità
# ============================================
slide7 = prs.slides[6]
set_title(slide7, 'Criticità e Opportunità')

add_subtitle(slide7, Emu(410714), Emu(1085649), Emu(5000000), Emu(400000), 'Criticità')

criticita = [
    '22 arbitri su 65 (34%) mai designati per il beach volley',
    '5 arbitri in attesa di rinnovo tesseramento',
    'Concentrazione su Romagna Uno: 18 arbitri (28% dell\'organico)',
    'Solo 25 arbitri attivi su 65 in organico (38%)',
]
add_bullet_textbox(slide7, Emu(410714), Emu(1500000), Emu(5500000), Emu(1800000), criticita, font_size=Pt(13))

add_subtitle(slide7, Emu(410714), Emu(3400000), Emu(5000000), Emu(400000), 'Opportunità')

opportunita = [
    'Coinvolgere comitati con pochi arbitri beach: Ferrara (2), Piacenza (3), Reggio Emilia (3)',
    'Attivare i 22 arbitri mai designati tramite formazione e affiancamento',
    'Trend positivo post-COVID: da 3 tornei (2020) a 16 (2024)',
    'Giovani promettenti: Avanzolini (27 anni, 15 tornei), Carota (19 anni, 10 tornei)',
]
add_bullet_textbox(slide7, Emu(410714), Emu(3850000), Emu(10500000), Emu(1800000), opportunita, font_size=Pt(13))

# Summary table on the right
add_subtitle(slide7, Emu(7000000), Emu(1085649), Emu(4500000), Emu(400000), 'Comitati da Potenziare')

pot_tbl = create_table(slide7, Emu(7000000), Emu(1500000), Emu(4200000), Emu(1500000), 5, 2)
for ci, h in enumerate(['Comitato', 'Arbitri Beach']):
    style_header_cell(pot_tbl.cell(0, ci), h)
for i, (ct, n) in enumerate([('Ferrara', '2'), ('Piacenza', '3'), ('Reggio Emilia', '3'), ('Parma', '5')]):
    fill = 'FFD9D9' if int(n) <= 3 else 'FFFFCC'
    style_data_cell(pot_tbl.cell(i+1, 0), ct, align=PP_ALIGN.LEFT, fill_color=fill)
    style_data_cell(pot_tbl.cell(i+1, 1), n, bold=True, fill_color=fill)


# ============================================
# SLIDE 8: Corso Arbitri Beach Volley (NEW)
# ============================================
slide8 = add_new_slide_with_base(prs)
set_title(slide8, 'Corso Arbitri Beach Volley')

add_subtitle(slide8, Emu(410714), Emu(1200000), Emu(5000000), Emu(400000), 'Dettagli del Corso')

placeholder_items = [
    'Date: in fase di definizione',
    'Sede: in fase di definizione',
    'Programma: in fase di definizione',
    'Docenti: in fase di definizione',
    'Destinatari: in fase di definizione',
]
add_bullet_textbox(slide8, Emu(410714), Emu(1700000), Emu(6000000), Emu(2500000), placeholder_items, font_size=Pt(14), color=RGBColor(0x66,0x66,0x66))

add_textbox(slide8, Emu(3000000), Emu(4500000), Emu(6000000), Emu(600000),
            '⚠ In fase di definizione',
            font_name='Calibri', font_size=Pt(24), bold=True, color=RGBColor(0xC0, 0x00, 0x00),
            alignment=PP_ALIGN.CENTER)


# ============================================
# SLIDE 9: Vuota template (backup)
# ============================================
slide9 = add_new_slide_with_base(prs)
set_title(slide9, 'Titolo presentazione')


# Save
prs.save('Conferenza_v2.pptx')
print('Saved Conferenza_v2.pptx successfully!')
print(f'Total slides: {len(prs.slides)}')
